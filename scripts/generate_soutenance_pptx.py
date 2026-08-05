"""Soutenance jalon final — 12 slides max, ton impersonnel, mise en page soignée."""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
CAP = ROOT / "captures"
OUT = ROOT / "soutenance" / "Jalon_Final_Soutenance_FSM_Docker.pptx"

BG = RGBColor(0x0A, 0x11, 0x1F)
CARD = RGBColor(0x12, 0x1A, 0x2B)
LINE = RGBColor(0x1E, 0x2A, 0x3F)
ACCENT = RGBColor(0x2D, 0xD4, 0xBF)
SKY = RGBColor(0x56, 0xCC, 0xF2)
TEXT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x8B, 0x9C, 0xB3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def font(p, size, bold=False, color=TEXT):
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"


def txt(p, text, size, bold=False, color=TEXT):
    p.text = text
    font(p, size, bold, color)


def box(slide, l, t, w, h, color, round_=True):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, l, t, w, h
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    if round_:
        sh.adjustments[0] = 0.06
    return sh


def bg(slide, prs):
    box(slide, 0, 0, prs.slide_width, prs.slide_height, BG, False)
    box(slide, 0, 0, Inches(0.12), prs.slide_height, ACCENT, False)


def footer(slide, i, n):
    b = slide.shapes.add_textbox(Inches(0.55), Inches(7.1), Inches(10), Inches(0.25))
    txt(b.text_frame.paragraphs[0], "FSM Docker  ·  Stage DevOps Wilance", 10, False, MUTED)
    b = slide.shapes.add_textbox(Inches(11.5), Inches(7.1), Inches(1.4), Inches(0.25))
    p = b.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    txt(p, f"{i} / {n}", 10, False, MUTED)


def heading(slide, title, subtitle=None):
    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.55))
    txt(t.text_frame.paragraphs[0], title, 28, True, WHITE)
    box(slide, Inches(0.55), Inches(0.88), Inches(1.5), Inches(0.05), ACCENT, False)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(1.0), Inches(12.2), Inches(0.35))
        txt(s.text_frame.paragraphs[0], subtitle, 13, False, MUTED)


def bullets(slide, items, l, t, w, h, size=15):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt(p, "•  " + item, size, False, TEXT)
        p.space_after = Pt(10)


def pick(*names: str) -> Path | None:
    for n in names:
        p = CAP / n
        if p.exists():
            return p
    return None


def fit(slide, path: Path | None, l, t, mw, mh):
    if not path:
        return
    use = path
    with Image.open(path) as im:
        w, h = im.size
        if path.suffix.lower() in {".jfif", ".jpeg", ".jpg", ".webp"}:
            tmp = path.with_name(path.stem + ".pptx_tmp.png")
            im.convert("RGB").save(tmp)
            use = tmp
    ratio = w / max(h, 1)
    if ratio > (mw / mh):
        width, height = mw, Emu(int(mw / ratio))
    else:
        height, width = mh, Emu(int(mh * ratio))
    x = l + (mw - width) // 2
    y = t + (mh - height) // 2
    slide.shapes.add_picture(str(use), x, y, width=width, height=height)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = []

    def add():
        s = prs.slides.add_slide(blank)
        slides.append(s)
        return s

    # ---- 1 Title ----
    s = add()
    bg(s, prs)
    box(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(3.9), CARD)
    t = s.shapes.add_textbox(Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.7))
    txt(t.text_frame.paragraphs[0], "Chaîne DevOps complète", 36, True, WHITE)
    t = s.shapes.add_textbox(Inches(0.9), Inches(2.8), Inches(11.5), Inches(0.5))
    txt(t.text_frame.paragraphs[0], "Application fil rouge FSM — du commit à l'observabilité", 20, False, ACCENT)
    t = s.shapes.add_textbox(Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.2))
    tf = t.text_frame
    txt(tf.paragraphs[0], "Stage Wilance  ·  Abdellah ANECLOUB & Youssef Ouchen", 16, False, TEXT)
    p = tf.add_paragraph()
    txt(p, "Git · CI/CD · Conteneurs · Cloud AWS · Kubernetes · GitOps · Monitoring", 14, False, MUTED)
    t = s.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.35))
    txt(t.text_frame.paragraphs[0], "fsm-app-morocco.duckdns.org", 13, False, SKY)

    # ---- 2 Contexte ----
    s = add()
    bg(s, prs)
    heading(s, "Contexte & objectif", "Une app simple pour construire une vraie chaîne")
    cards = [
        ("Application", "Next.js + PostgreSQL\nGestion d'interventions FSM"),
        ("Objectif", "Automatiser, sécuriser,\ndéployer et observer"),
        ("Méthode", "Jalons successifs,\nPR, preuves, démos"),
    ]
    for i, (h, d) in enumerate(cards):
        left = Inches(0.55) + i * Inches(4.15)
        box(s, left, Inches(1.7), Inches(3.95), Inches(4.6), CARD)
        tb = s.shapes.add_textbox(left + Inches(0.3), Inches(2.2), Inches(3.35), Inches(0.5))
        txt(tb.text_frame.paragraphs[0], h, 18, True, ACCENT)
        tb = s.shapes.add_textbox(left + Inches(0.3), Inches(3.0), Inches(3.35), Inches(2.5))
        p = tb.text_frame.paragraphs[0]
        p.word_wrap = True
        txt(p, d, 16, False, TEXT)

    # ---- 3 Architecture ----
    s = add()
    bg(s, prs)
    heading(s, "Architecture de la chaîne", "Flux unique de bout en bout")
    steps = [
        ("GitHub", "PR & revue"),
        ("CI/CD", "Tests & sécu"),
        ("GHCR", "Image Docker"),
        ("AWS EC2", "Terraform\nAnsible"),
        ("Kubernetes", "Helm\nArgo CD"),
        ("Observabilité", "Prometheus\nGrafana"),
    ]
    for i, (h, d) in enumerate(steps):
        left = Inches(0.4) + i * Inches(2.15)
        box(s, left, Inches(1.8), Inches(2.0), Inches(3.4), CARD)
        tb = s.shapes.add_textbox(left + Inches(0.1), Inches(2.15), Inches(1.8), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        txt(p, str(i + 1), 20, True, ACCENT)
        tb = s.shapes.add_textbox(left + Inches(0.1), Inches(2.9), Inches(1.8), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        txt(p, h, 14, True, WHITE)
        tb = s.shapes.add_textbox(left + Inches(0.1), Inches(3.55), Inches(1.8), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        txt(p, d, 12, False, MUTED)
        if i < len(steps) - 1:
            a = s.shapes.add_textbox(left + Inches(1.85), Inches(3.1), Inches(0.35), Inches(0.4))
            txt(a.text_frame.paragraphs[0], "→", 16, True, SKY)
    box(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(1.2), CARD)
    bullets(
        s,
        [
            "Prod web HTTPS sur EC2 (DuckDNS) + orchestration K8s/GitOps en parallèle",
            "Même image GHCR utilisée sur les cibles de déploiement",
        ],
        Inches(0.8),
        Inches(5.7),
        Inches(11.7),
        Inches(0.9),
        14,
    )

    # ---- 4 Fondations + Docker ----
    s = add()
    bg(s, prs)
    heading(s, "Fondations & conteneurisation", "Git propre · Docker multi-stage · Compose · Postgres local")
    box(s, Inches(0.55), Inches(1.6), Inches(5.9), Inches(5.0), CARD)
    bullets(
        s,
        [
            "Travail par branches et Pull Requests",
            "Route de santé /api/health",
            "Image optimisée (~264 Mo)",
            "Stack locale : app + PostgreSQL",
            "Auth locale (NextAuth) découplée du SaaS",
            "Données persistantes via volume",
        ],
        Inches(0.85),
        Inches(1.95),
        Inches(5.4),
        Inches(4.4),
        15,
    )
    box(s, Inches(6.7), Inches(1.6), Inches(6.1), Inches(5.0), CARD)
    fit(s, pick("dbeaver.png", "navigateur.png"), Inches(6.9), Inches(1.85), Inches(5.7), Inches(4.5))

    # ---- 5 CI/CD + DevSecOps ----
    s = add()
    bg(s, prs)
    heading(s, "CI/CD & sécurité", "Automatiser la qualité avant la publication")
    box(s, Inches(0.55), Inches(1.55), Inches(4.5), Inches(5.1), CARD)
    bullets(
        s,
        [
            "GitHub Actions sur PR / main",
            "Tests unitaires",
            "Gitleaks (secrets)",
            "Trivy (vulnérabilités)",
            "Build & push GHCR",
            "Pipeline bloquant si échec",
        ],
        Inches(0.8),
        Inches(1.9),
        Inches(4.0),
        Inches(4.5),
        15,
    )
    box(s, Inches(5.3), Inches(1.55), Inches(7.5), Inches(5.1), CARD)
    fit(s, pick("jalon-final-ci-green.png", "ci-success.png", "jalon5-pipline-vert.png"), Inches(5.5), Inches(1.8), Inches(7.1), Inches(4.6))

    # ---- 6 Deploy cloud ----
    s = add()
    bg(s, prs)
    heading(s, "Déploiement cloud", "Automatisation jusqu'à la mise en ligne")
    box(s, Inches(0.55), Inches(1.55), Inches(5.7), Inches(5.1), CARD)
    bullets(
        s,
        [
            "Base managée Neon",
            "Déploiement auto (Railway) via CI",
            "Health check après mise en ligne",
            "Procédure de rollback documentée",
            "Secrets hors dépôt (GitHub Secrets)",
        ],
        Inches(0.85),
        Inches(1.95),
        Inches(5.2),
        Inches(4.4),
        15,
    )
    box(s, Inches(6.5), Inches(1.55), Inches(6.3), Inches(5.1), CARD)
    fit(s, pick("jalon-final-site-duckdns.png", "railway_deploy.jfif"), Inches(6.7), Inches(1.8), Inches(5.9), Inches(4.6))

    # ---- 7 IaC AWS ----
    s = add()
    bg(s, prs)
    heading(s, "Infrastructure as Code", "AWS EC2 reproductible avec Terraform & Ansible")
    box(s, Inches(0.55), Inches(1.55), Inches(5.7), Inches(5.1), CARD)
    bullets(
        s,
        [
            "Terraform : EC2 + Security Group",
            "Ansible : Docker, Nginx, Certbot",
            "DuckDNS + HTTPS Let's Encrypt",
            "Infra séparée (dépôt dédié)",
            "Destroy / recreate démontré",
            "Idempotence du playbook",
        ],
        Inches(0.85),
        Inches(1.95),
        Inches(5.2),
        Inches(4.4),
        15,
    )
    box(s, Inches(6.5), Inches(1.55), Inches(6.3), Inches(5.1), CARD)
    fit(s, pick("jalon-final-site-duckdns.png"), Inches(6.7), Inches(1.8), Inches(5.9), Inches(4.6))
    cap = s.shapes.add_textbox(Inches(6.7), Inches(6.15), Inches(5.9), Inches(0.3))
    p = cap.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    txt(p, "https://fsm-app-morocco.duckdns.org", 12, False, SKY)

    # ---- 8 K8s + Helm ----
    s = add()
    bg(s, prs)
    heading(s, "Kubernetes & Helm", "Orchestration et paquet réutilisable")
    box(s, Inches(0.55), Inches(1.55), Inches(5.5), Inches(5.1), CARD)
    bullets(
        s,
        [
            "Deployment, Service, Ingress",
            "Mise à l'échelle des réplicas",
            "Sondes /api/health + limites ressources",
            "Chart Helm fsm-app + values.yaml",
            "Upgrade et rollback validés",
        ],
        Inches(0.85),
        Inches(1.95),
        Inches(5.0),
        Inches(4.4),
        15,
    )
    box(s, Inches(6.3), Inches(1.55), Inches(6.5), Inches(2.2), CARD)
    fit(s, pick("jalon-final-pods-running.png", "jalon7-kubectl-get-pods.png"), Inches(6.5), Inches(1.7), Inches(6.1), Inches(1.9))
    box(s, Inches(6.3), Inches(3.95), Inches(6.5), Inches(2.7), CARD)
    fit(s, pick("preuve-helm.jfif", "jalon8-helm-template-test-local.png", "jalon7-app-validation.png"), Inches(6.5), Inches(4.1), Inches(6.1), Inches(2.4))

    # ---- 9 GitOps ----
    s = add()
    bg(s, prs)
    heading(s, "GitOps avec Argo CD", "Git comme unique source de vérité")
    box(s, Inches(0.55), Inches(1.55), Inches(7.6), Inches(5.1), CARD)
    fit(s, pick("jalon9-argocd-healthy-synced.png", "Argo_CD_Synced.jfif"), Inches(0.75), Inches(1.75), Inches(7.2), Inches(4.7))
    box(s, Inches(8.4), Inches(1.55), Inches(4.4), Inches(5.1), CARD)
    bullets(
        s,
        [
            "Sync automatique",
            "Self-heal activé",
            "Prune des ressources",
            "État Healthy / Synced",
            "Drift corrigé depuis Git",
        ],
        Inches(8.7),
        Inches(2.1),
        Inches(3.9),
        Inches(4.2),
        15,
    )

    # ---- 10 Observabilité ----
    s = add()
    bg(s, prs)
    heading(s, "Observabilité & fiabilité", "Mesurer, visualiser, alerter")
    box(s, Inches(0.55), Inches(1.5), Inches(4.3), Inches(5.15), CARD)
    bullets(
        s,
        [
            "Métriques /api/metrics",
            "Collecte Prometheus",
            "Dashboard Grafana",
            "SLO : p95 < 1s",
            "Alerte haute latence",
        ],
        Inches(0.85),
        Inches(1.9),
        Inches(3.8),
        Inches(4.5),
        15,
    )
    box(s, Inches(5.1), Inches(1.5), Inches(7.7), Inches(3.15), CARD)
    fit(s, pick("jalon10-grafana-dashboard.png"), Inches(5.3), Inches(1.65), Inches(7.3), Inches(2.85))
    box(s, Inches(5.1), Inches(4.85), Inches(7.7), Inches(1.8), CARD)
    fit(s, pick("jalon10-alert-fsm-high-latency-firing.png"), Inches(5.3), Inches(4.95), Inches(7.3), Inches(1.6))

    # ---- 11 Démo / résultats ----
    s = add()
    bg(s, prs)
    heading(s, "Résultat démontrable", "La chaîne fonctionne de bout en bout")
    items = [
        ("CI verte", "Qualité + sécu automatisées"),
        ("Site live", "HTTPS DuckDNS"),
        ("Cluster", "Pods Ready / Helm"),
        ("GitOps", "Argo Synced"),
        ("Monitoring", "Grafana + alerte"),
        ("Repo propre", "Secrets hors Git"),
    ]
    for i, (h, d) in enumerate(items):
        col, row = i % 3, i // 3
        left = Inches(0.55) + col * Inches(4.15)
        top = Inches(1.65) + row * Inches(2.5)
        box(s, left, top, Inches(3.95), Inches(2.25), CARD)
        tb = s.shapes.add_textbox(left + Inches(0.3), top + Inches(0.45), Inches(3.35), Inches(0.5))
        txt(tb.text_frame.paragraphs[0], h, 18, True, ACCENT)
        tb = s.shapes.add_textbox(left + Inches(0.3), top + Inches(1.15), Inches(3.35), Inches(0.6))
        txt(tb.text_frame.paragraphs[0], d, 14, False, TEXT)

    # ---- 12 Clôture ----
    s = add()
    bg(s, prs)
    heading(s, "Apprentissages & suite", "Merci — Questions ?")
    box(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(5.0), CARD)
    bullets(
        s,
        [
            "Automatiser tôt (CI + sécu)",
            "Déclarer l'infra et les déploiements",
            "Git comme source de vérité",
            "Observer pour fiabiliser",
            "Prouver chaque étape",
        ],
        Inches(0.9),
        Inches(2.0),
        Inches(5.4),
        Inches(4.3),
        16,
    )
    box(s, Inches(6.85), Inches(1.6), Inches(5.95), Inches(5.0), CARD)
    b = s.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.3), Inches(4.0))
    tf = b.text_frame
    tf.word_wrap = True
    txt(tf.paragraphs[0], "Perspectives", 16, True, ACCENT)
    for line in [
        "",
        "• Industrialiser les secrets (Sealed Secrets)",
        "• Stabiliser les réplicas",
        "• Renforcer les dashboards / SLO",
        "",
        "Repos",
        "github.com/abdellah-get/fsm-docker",
        "fsm-app-morocco.duckdns.org",
    ]:
        p = tf.add_paragraph()
        txt(p, line, 14, False, TEXT)

    n = len(slides)
    assert n <= 12, n
    for i, slide in enumerate(slides, 1):
        footer(slide, i, n)

    for tmp in CAP.glob("*.pptx_tmp.png"):
        tmp.unlink(missing_ok=True)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"{n} slides → {OUT}")


if __name__ == "__main__":
    build()
